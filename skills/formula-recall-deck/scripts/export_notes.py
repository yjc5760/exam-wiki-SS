#!/usr/bin/env python3
"""把 .pptx 每頁的備忘稿匯出成逐頁旁白稿 Markdown。

自動偵測回想卡頁（recallSlide 產生的「回想卡 N / M」標記）、統計旁白字數、
估算配音片長，並在旁白稿開頭列出需要插入 3 秒靜音的頁碼。

執行：python3 export_notes.py "XX-Un-m_主題_公式給背分界_記憶片.pptx" [輸出檔名.md]
"""
import re
import sys
from pathlib import Path
from pptx import Presentation

CPS = 4.5  # 中文語速：字／秒


def slide_texts(slide):
    return [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]


def biggest_text(slide):
    """取投影片上字級最大的一段文字當作該頁標題。"""
    title, best = "", -1
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                size = run.font.size.pt if run.font.size else 0
                if size > best and run.text.strip():
                    best, title = size, run.text.strip().replace("\n", " ")
    return title


def main():
    if len(sys.argv) < 2:
        sys.exit("用法：python3 export_notes.py <deck.pptx> [out.md]")
    pptx_path = Path(sys.argv[1])
    prs = Presentation(pptx_path)
    slides = list(prs.slides)

    stem = pptx_path.stem
    code = stem.split("_")[0]  # 例：SS-U1-1
    out_path = Path(sys.argv[2]) if len(sys.argv) > 2 else Path(f"旁白稿_{code}.md")

    recall_pages = [
        i for i, s in enumerate(slides, 1)
        if any(re.match(r"^回想卡 \d+ */ *\d+$", t.strip()) for t in slide_texts(s))
    ]

    body, missing, total_chars = [], [], 0
    for i, s in enumerate(slides, 1):
        note = s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
        total_chars += len(note)
        if not note:
            missing.append(i)
        tag = "　⏸ **回想卡**" if i in recall_pages else ""
        body += [f"## P{i:02d}　{biggest_text(s)}{tag}", "", note or "_（無旁白）_"]
        if i in recall_pages:
            body += ["", "`旁白結束後插入 3 秒靜音（讓觀眾默寫）`"]
        body += [""]

    secs = total_chars / CPS
    lo, hi = int(secs / 60), int(secs / 60) + 2
    pages = "、".join(f"P{n:02d}" for n in recall_pages) or "（無）"

    head = [
        f"# {code}｜逐頁旁白稿",
        "",
        f"對應簡報：`{pptx_path.with_suffix('.pdf').name}`（共 {len(slides)} 頁）  ",
        "用途：餵給 `pdf-narration-video` pipeline 做 Azure TTS 配音＋字幕對齊。",
        "",
        "**製作備註**",
        "",
        f"- 全片旁白約 {total_chars:,} 字；中文語速取 {CPS} 字／秒，配音後預估 {lo}–{hi} 分鐘。",
        f"- 回想卡頁（共 {len(recall_pages)} 頁：{pages}）請在旁白結束後**額外插入 3 秒靜音**，讓觀眾有時間默寫。",
        "- 旁白中的公式已改寫成口語唸法，可直接送 TTS，不需再處理數學符號。",
        "- 同一份旁白也寫進 .pptx 的「備忘稿」欄位，用 PowerPoint 自行錄製時可直接對照。",
        "",
        "---",
        "",
    ]

    out_path.write_text("\n".join(head + body), encoding="utf-8")
    print(f"已輸出 {out_path}")
    print(f"  頁數 {len(slides)}　旁白 {total_chars:,} 字　預估 {lo}–{hi} 分鐘")
    print(f"  回想卡頁：{pages}")
    if missing:
        print(f"  ⚠ 這些頁沒有旁白，請補：{missing}")


if __name__ == "__main__":
    main()
